from pptx import Presentation
from pathlib import Path

OUTPUT_DIR = Path("docs")
OUTPUT_DIR.mkdir(exist_ok=True)


def create_slide_deck():
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Customer Churn Prediction"
    title_slide.placeholders[1].text = "ML Data Associate Portfolio Project"

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Project Objective"
    slide1.placeholders[1].text = (
        "- Predict customer churn using historical customer data\n"
        "- Help business teams reduce customer churn\n"
        "- Improve customer retention strategies and revenue outcomes"
    )

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Findings"
    slide2.placeholders[1].text = (
        "- Shorter customer tenure is linked to higher churn risk\n"
        "- Month-to-month contracts increase cancellation risk\n"
        "- Higher monthly charges and weak support signals increase churn\n"
        "- Model delivers strong predictive performance for retention planning"
    )

    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Methodology"
    slide3.placeholders[1].text = (
        "- Data cleaning and preprocessing\n"
        "- Exploratory data analysis\n"
        "- Feature engineering and model comparison\n"
        "- Evaluation with accuracy, precision, recall, F1, ROC-AUC\n"
        "- Deployment through Streamlit dashboard"
    )

    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Business Impact"
    slide4.placeholders[1].text = (
        "- Reduce revenue loss from churn\n"
        "- Enable proactive retention campaigns\n"
        "- Prioritize support for high-risk customers\n"
        "- Improve customer lifetime value"
    )

    output_file = OUTPUT_DIR / "customer_churn_presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved to: {output_file}")


if __name__ == "__main__":
    create_slide_deck()
